"""Slide preview, validation, and layout templates for PPTX generation.

Usage from Claude Code:
    python3 /Users/ruichenzhao/projects/twpa/.claude/slide_tools.py preview PPTX_PATH [SLIDE_NUM]
    python3 /Users/ruichenzhao/projects/twpa/.claude/slide_tools.py validate PPTX_PATH [SLIDE_NUM]
    python3 /Users/ruichenzhao/projects/twpa/.claude/slide_tools.py preview-all PPTX_PATH
"""
import sys, io, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

# ─── Constants ───────────────────────────────────────────────────────────────
SLIDE_W_IN = 13.333  # 16:9 widescreen
SLIDE_H_IN = 7.5
PREVIEW_DPI = 150
MIN_GAP_IN = 0.15     # minimum gap between shapes
MIN_MARGIN_IN = 0.25  # minimum margin from slide edge

# Font size rules (projection-readable minimums)
FONT_RULES = {
    "matplotlib": {"axis_label": 14, "title": 15, "tick_label": 12,
                   "annotation": 10, "legend": 11, "formula": 18},
    "plotly":     {"main_font": 20, "title": 24, "tick_font": 20,
                   "legend": 16, "annotation": 14},
    "pptx":       {"slide_title": 22, "subtitle": 14, "body": 12,
                   "caption": 10, "footer": 9},
}


def emu_to_in(emu):
    return emu / 914400


def shape_bbox(shape):
    """Return (left, top, right, bottom) in inches."""
    l = emu_to_in(shape.left)
    t = emu_to_in(shape.top)
    w = emu_to_in(shape.width)
    h = emu_to_in(shape.height)
    return (l, t, l + w, t + h)


# ─── Preview ─────────────────────────────────────────────────────────────────

def preview_slide(pptx_path, slide_num, dpi=PREVIEW_DPI):
    """Render a composite preview of a single slide. Returns PIL Image."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num - 1]  # 1-indexed

    sw = emu_to_in(prs.slide_width)
    sh = emu_to_in(prs.slide_height)
    canvas = Image.new('RGB', (int(sw * dpi), int(sh * dpi)), '#0a0a0a')
    draw = ImageDraw.Draw(canvas)

    try:
        font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 14)
        font_small = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 11)
    except Exception:
        font = ImageFont.load_default()
        font_small = font

    for shape in slide.shapes:
        x = int(emu_to_in(shape.left) * dpi)
        y = int(emu_to_in(shape.top) * dpi)
        w = int(emu_to_in(shape.width) * dpi)
        h = int(emu_to_in(shape.height) * dpi)

        if w <= 0 or h <= 0:
            continue

        if shape.shape_type == 13:  # PICTURE
            try:
                img_data = shape.image.blob
                img = Image.open(io.BytesIO(img_data))
                img = img.resize((w, h), Image.LANCZOS)
                canvas.paste(img, (x, y))
            except Exception:
                draw.rectangle([x, y, x + w, y + h], outline='red', width=2)
                draw.text((x + 4, y + 4), '[IMAGE]', fill='red', font=font)
        elif shape.has_text_frame:
            text = shape.text_frame.text[:120]
            draw.rectangle([x, y, x + w, y + h], outline='#444444', width=1)
            draw.text((x + 4, y + 4), text, fill='#f0f0f0', font=font_small)

    # Draw slide border
    draw.rectangle([0, 0, int(sw * dpi) - 1, int(sh * dpi) - 1],
                   outline='#555555', width=2)

    return canvas


def preview_all(pptx_path, dpi=100):
    """Preview all slides in a grid."""
    prs = Presentation(pptx_path)
    n = len(prs.slides)
    cols = 4
    rows = (n + cols - 1) // cols

    sw = emu_to_in(prs.slide_width)
    sh = emu_to_in(prs.slide_height)
    thumb_w = int(sw * dpi)
    thumb_h = int(sh * dpi)
    pad = 10

    grid = Image.new('RGB',
                     (cols * (thumb_w + pad) + pad, rows * (thumb_h + pad + 20) + pad),
                     '#1a1a1a')
    draw = ImageDraw.Draw(grid)

    try:
        font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 14)
    except Exception:
        font = ImageFont.load_default()

    for i in range(n):
        col = i % cols
        row = i // cols
        x = pad + col * (thumb_w + pad)
        y = pad + row * (thumb_h + pad + 20)

        thumb = preview_slide(pptx_path, i + 1, dpi=dpi)
        grid.paste(thumb, (x, y + 18))
        draw.text((x + 4, y + 2), f'Slide {i + 1}', fill='#ffd60a', font=font)

    return grid


# ─── Validation ──────────────────────────────────────────────────────────────

def validate_slide(pptx_path, slide_num):
    """Check a slide against layout rules. Returns list of issues."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num - 1]
    issues = []

    sw = emu_to_in(prs.slide_width)
    sh = emu_to_in(prs.slide_height)
    shapes_info = []

    for shape in slide.shapes:
        bbox = shape_bbox(shape)
        l, t, r, b = bbox
        w = r - l
        h = b - t
        is_pic = shape.shape_type == 13
        is_text = shape.has_text_frame
        text = ""
        if is_text:
            text = shape.text_frame.text[:60]

        shapes_info.append({
            'name': shape.name, 'bbox': bbox, 'is_pic': is_pic,
            'is_text': is_text, 'text': text, 'w': w, 'h': h,
        })

        # Check 1: overflow beyond slide
        if r > sw + 0.1:
            issues.append(f'OVERFLOW: {shape.name} right edge ({r:.2f}") exceeds slide width ({sw:.2f}")')
        if b > sh + 0.1:
            issues.append(f'OVERFLOW: {shape.name} bottom ({b:.2f}") exceeds slide height ({sh:.2f}")')

        # Check 2: text font sizes and contrast (for PPTX textboxes)
        if is_text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if run.font.size:
                        pts = run.font.size / 12700  # EMU to pt
                        if pts < 9:
                            issues.append(f'TINY TEXT: {shape.name} has {pts:.0f}pt text ("{run.text[:30]}") — min 9pt')
                    # Check contrast: dark/unset text on dark background is invisible
                    try:
                        rgb = run.font.color.rgb
                    except AttributeError:
                        rgb = None
                    if rgb is None:
                        issues.append(
                            f'INVISIBLE TEXT: {shape.name} has inherited/unset color '
                            f'("{run.text[:40]}") — will be black on dark bg. Set explicitly.')
                    elif rgb[0] < 0x30 and rgb[1] < 0x30 and rgb[2] < 0x30:
                        issues.append(
                            f'INVISIBLE TEXT: {shape.name} has near-black color #{rgb} '
                            f'("{run.text[:40]}") — invisible on dark bg.')

    # Check 3: pairwise overlap between pictures and text
    for i, a in enumerate(shapes_info):
        for j, b_info in enumerate(shapes_info):
            if j <= i:
                continue
            al, at, ar, ab = a['bbox']
            bl, bt, br, bb = b_info['bbox']
            # Check overlap
            if al < br and ar > bl and at < bb and ab > bt:
                # Compute overlap area
                ox = min(ar, br) - max(al, bl)
                oy = min(ab, bb) - max(at, bt)
                area = ox * oy
                # Skip tiny text-text overlaps (intentional bullet stacking)
                if not (a['is_pic'] or b_info['is_pic']) and area < 0.15:
                    continue
                if area > 0.05:  # significant overlap (>0.05 sq in)
                    issues.append(
                        f'OVERLAP: {a["name"]} and {b_info["name"]} overlap by {ox:.2f}"×{oy:.2f}" '
                        f'({area:.2f} sq in)')

    # Check 4: pictures too close together
    pics = [s for s in shapes_info if s['is_pic']]
    for i, a in enumerate(pics):
        for j, b_info in enumerate(pics):
            if j <= i:
                continue
            al, at, ar, ab = a['bbox']
            bl, bt, br, bb = b_info['bbox']
            # Vertical gap
            if al < br and ar > bl:  # horizontally overlapping
                gap = bt - ab if bt > ab else at - bb
                if 0 < gap < MIN_GAP_IN:
                    issues.append(
                        f'CRAMPED: {a["name"]} and {b_info["name"]} only {gap:.2f}" apart vertically (min {MIN_GAP_IN}")')
            # Horizontal gap
            if at < bb and ab > bt:  # vertically overlapping
                gap = bl - ar if bl > ar else al - br
                if 0 < gap < MIN_GAP_IN:
                    issues.append(
                        f'CRAMPED: {a["name"]} and {b_info["name"]} only {gap:.2f}" apart horizontally (min {MIN_GAP_IN}")')

    # Check 5: space utilization
    content_area = 0
    for s in shapes_info:
        content_area += s['w'] * s['h']
    slide_area = sw * sh
    util = content_area / slide_area
    if util < 0.4:
        issues.append(f'SPARSE: content covers only {util:.0%} of slide area (target ≥50%)')

    return issues


def validate_all(pptx_path):
    """Validate all slides."""
    prs = Presentation(pptx_path)
    all_issues = {}
    for i in range(len(prs.slides)):
        issues = validate_slide(pptx_path, i + 1)
        if issues:
            all_issues[i + 1] = issues
    return all_issues


# ─── Layout Templates ────────────────────────────────────────────────────────

class SlideLayout:
    """Pre-calculated layout slots for common slide patterns on 16:9 slides.

    All positions in inches. Assumes:
    - Title area: T=0.20 to B=1.00
    - Footer: T=7.10 to B=7.45
    - Content area: T=1.05 to B=7.00
    """
    TITLE_TOP = 0.20
    TITLE_H = 0.70
    CONTENT_TOP = 1.05
    CONTENT_BOTTOM = 7.00
    FOOTER_TOP = 7.10
    MARGIN_L = 0.35
    MARGIN_R = 0.35
    CONTENT_W = SLIDE_W_IN - MARGIN_L - MARGIN_R  # ~12.63"

    @classmethod
    def single_image(cls, image_ar, with_bullets=False):
        """Full-width single image, optional bullet text on right.
        Returns dict with 'image' and optionally 'bullets' slot positions."""
        content_h = cls.CONTENT_BOTTOM - cls.CONTENT_TOP

        if with_bullets:
            bullet_w = 2.80
            img_w = cls.CONTENT_W - bullet_w - 0.20  # gap
            img_h = min(content_h, img_w / image_ar)
            img_w = img_h * image_ar  # recalc for aspect ratio
            img_t = cls.CONTENT_TOP + (content_h - img_h) / 2
            return {
                'image': (cls.MARGIN_L, img_t, img_w, img_h),
                'bullets': (cls.MARGIN_L + img_w + 0.20, cls.CONTENT_TOP, bullet_w, content_h),
            }
        else:
            img_w = cls.CONTENT_W
            img_h = min(content_h, img_w / image_ar)
            img_w = img_h * image_ar
            img_l = cls.MARGIN_L + (cls.CONTENT_W - img_w) / 2
            img_t = cls.CONTENT_TOP + (content_h - img_h) / 2
            return {
                'image': (img_l, img_t, img_w, img_h),
            }

    @classmethod
    def two_column(cls, ar_left, ar_right, gap=0.30):
        """Two images side by side.
        Returns dict with 'left' and 'right' slot positions."""
        content_h = cls.CONTENT_BOTTOM - cls.CONTENT_TOP
        half_w = (cls.CONTENT_W - gap) / 2

        lh = min(content_h, half_w / ar_left)
        lw = lh * ar_left
        rh = min(content_h, half_w / ar_right)
        rw = rh * ar_right

        lt = cls.CONTENT_TOP + (content_h - lh) / 2
        rt = cls.CONTENT_TOP + (content_h - rh) / 2
        ll = cls.MARGIN_L + (half_w - lw) / 2
        rl = cls.MARGIN_L + half_w + gap + (half_w - rw) / 2

        return {
            'left': (ll, lt, lw, lh),
            'right': (rl, rt, rw, rh),
        }

    @classmethod
    def stacked(cls, ar_top, ar_bottom, gap=0.20, caption_h=0.25):
        """Two images stacked vertically with optional captions.
        Returns dict with 'top_caption', 'top', 'bottom_caption', 'bottom'."""
        avail_h = cls.CONTENT_BOTTOM - cls.CONTENT_TOP - 2 * caption_h - gap
        each_h = avail_h / 2

        top_w = min(cls.CONTENT_W, each_h * ar_top)
        top_h = top_w / ar_top
        bot_w = min(cls.CONTENT_W, each_h * ar_bottom)
        bot_h = bot_w / ar_bottom

        cap1_t = cls.CONTENT_TOP
        top_t = cap1_t + caption_h
        cap2_t = top_t + top_h + gap
        bot_t = cap2_t + caption_h

        return {
            'top_caption': (cls.MARGIN_L, cap1_t, cls.CONTENT_W, caption_h),
            'top': (cls.MARGIN_L + (cls.CONTENT_W - top_w) / 2, top_t, top_w, top_h),
            'bottom_caption': (cls.MARGIN_L, cap2_t, cls.CONTENT_W, caption_h),
            'bottom': (cls.MARGIN_L + (cls.CONTENT_W - bot_w) / 2, bot_t, bot_w, bot_h),
        }


# ─── CLI ─────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("Usage: slide_tools.py <preview|validate|preview-all> PPTX_PATH [SLIDE_NUM]")
        sys.exit(1)

    cmd = sys.argv[1]
    pptx_path = sys.argv[2]
    slide_num = int(sys.argv[3]) if len(sys.argv) > 3 else None

    if cmd == 'preview':
        if not slide_num:
            print("Slide number required for preview")
            sys.exit(1)
        img = preview_slide(pptx_path, slide_num)
        out = f'/tmp/slide{slide_num}_preview.png'
        img.save(out)
        print(f'Preview saved: {out}')

    elif cmd == 'preview-all':
        img = preview_all(pptx_path)
        out = '/tmp/deck_preview.png'
        img.save(out)
        print(f'Full deck preview saved: {out}')

    elif cmd == 'validate':
        if slide_num:
            issues = validate_slide(pptx_path, slide_num)
            if issues:
                print(f'Slide {slide_num}: {len(issues)} issues')
                for iss in issues:
                    print(f'  ❌ {iss}')
            else:
                print(f'Slide {slide_num}: ✓ no issues')
        else:
            all_issues = validate_all(pptx_path)
            if all_issues:
                for sn, issues in sorted(all_issues.items()):
                    print(f'Slide {sn}: {len(issues)} issues')
                    for iss in issues:
                        print(f'  ❌ {iss}')
            else:
                print('All slides: ✓ no issues')
            total = sum(len(v) for v in all_issues.values())
            print(f'\nTotal: {total} issues across {len(all_issues)} slides')

    else:
        print(f"Unknown command: {cmd}")
        sys.exit(1)


if __name__ == '__main__':
    main()
